import os
import csv
import pypdf
import docx
import openpyxl
import pptx
from typing import Optional
from gemma4.manager import get_manager

def read_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf(file_path: str) -> str:
    text = ""
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
    return text

def read_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def read_pptx(file_path: str) -> str:
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        text += f"Sheet: {sheet}\n"
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

def read_csv(file_path: str) -> str:
    text = ""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            text += "\t".join(row) + "\n"
    return text

def read_file_content(file_path: str) -> str:
    """
    Tự động nhận diện định dạng và trích xuất nội dung văn bản từ file.
    """
    if not os.path.exists(file_path):
        return f"Lỗi: Không tìm thấy file tại {file_path}"
    
    ext = os.path.splitext(file_path)[1].lower()

    # Kiểm tra header cho file không có extension hoặc extension lạ
    try:
        if os.path.exists(file_path):
            with open(file_path, "rb") as f:
                header = f.read(64)
            if header:
                # Audio magic bytes (RIFF WAVE, ID3, OggS, fLaC, mp3)
                if (header.startswith(b"RIFF") and b"WAVE" in header[:16]) or header.startswith(b"ID3") or header.startswith(b"OggS") or header.startswith(b"fLaC") or header.startswith(b"\xff\xfb") or header.startswith(b"\xff\xf3"):
                    from myassitant.file_worker import _transcribe_audio_local
                    desc = _transcribe_audio_local(file_path)
                    return desc if desc else f"[File âm thanh: {os.path.basename(file_path)}]"
                # Image magic bytes
                if header.startswith(b"\xff\xd8\xff") or header.startswith(b"\x89PNG") or header.startswith(b"GIF8"):
                    from myassitant.file_worker import _describe_image_via_api
                    desc = _describe_image_via_api(file_path)
                    return desc if desc else f"[File ảnh: {os.path.basename(file_path)}]"
    except Exception:
        pass
    
    try:
        if ext in [".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"]:
            try:
                from myassitant.file_worker import _describe_image_via_api
                desc = _describe_image_via_api(file_path)
                return desc if desc else f"[File ảnh: {os.path.basename(file_path)}]"
            except Exception as e:
                return f"[File ảnh {os.path.basename(file_path)} - Không mô tả được: {e}]"
        elif ext in [".mp4", ".mkv", ".avi", ".mov", ".webm", ".flv", ".3gp", ".ts", ".mp3", ".wav", ".flac", ".ogg", ".opus", ".m4a", ".aac", ".oga"]:
            try:
                from myassitant.file_worker import _transcribe_audio_local
                desc = _transcribe_audio_local(file_path)
                return desc if desc else f"[File media {os.path.basename(file_path)}]"
            except Exception as e:
                return f"[File media {os.path.basename(file_path)} - Không xử lý được: {e}]"
        elif ext in [".txt", ".log", ".md", ".json", ".py"]:
            return read_txt(file_path)
        elif ext == ".pdf":
            return read_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return read_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            return read_pptx(file_path)
        elif ext in [".xlsx", ".xls"]:
            return read_xlsx(file_path)
        elif ext == ".csv":
            return read_csv(file_path)
        else:
            # Fallback cho định dạng khác — kiểm tra xem có phải binary không
            try:
                with open(file_path, "rb") as f:
                    sample = f.read(512)
                text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7f})
                is_binary = bool(sample.translate(None, text_chars))
                if is_binary:
                    return f"[File nhị phân {os.path.basename(file_path)} không thể đọc dạng văn bản]"
            except Exception:
                pass
            return read_txt(file_path)
    except Exception as e:
        return f"Lỗi khi đọc file {ext}: {str(e)}"

def process_file_with_prompt(file_path: str, prompt: str, model_id: str = "unsloth/gemma-4-e4b-it-unsloth-bnb-4bit") -> str:
    """
    Đọc nội dung file và đưa vào Gemma 4 kèm theo prompt để xử lý.
    """
    content = read_file_content(file_path)
    if content.startswith("Lỗi:"):
        return content
    
    # Truncate content to avoid exceeding context limits (approx 4k tokens)
    # 1 token ~ 4 characters for English, ~2.5 characters for Vietnamese. 
    # Use 12000 chars as a conservative limit.
    max_chars = 12000
    if len(content) > max_chars:
        content = content[:max_chars] + "\n... [Nội dung đã được rút gọn do quá dài] ..."
        
    full_prompt = f"Nội dung của file:\n{content}\n\nYêu cầu: {prompt}"
    
    manager = get_manager(model_id)
    return manager.generate(full_prompt)
