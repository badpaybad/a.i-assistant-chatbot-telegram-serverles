import os
import datetime
import pickle
import json
import numpy as np
import faiss
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from google import genai
from google.genai import types
from config import TELEGRAM_BOT_TOKEN, TELEGRAM_API_URL, PORT, TELEGRAM_BOT_CHATID, TELEGRAM_BOT_USERNAME, GEMINI_APIKEY, DISCORD_PUBKEY, DISCORD_APPID, DISCORD_TOKEN,  TELEGRAM_API_ID, TELEGRAM_API_HASH, REPLY_ON_TAG_BOT_USERNAME

class LocalVectorDB:
    def __init__(self, db_name,embedding_model_type="models/gemini-embedding-001"):
        self.db_name = db_name
        self.db_path = f"data/vector_db/{db_name}"
        self.index_file = f"{self.db_path}/index.faiss"
        self.metadata_file = f"{self.db_path}/metadata.pkl"
        self.client = genai.Client(api_key=GEMINI_APIKEY)
        self.embedding_model_type=embedding_model_type
        if self.embedding_model_type=="models/gemini-embedding-001":
            self.embedding_model = "models/gemini-embedding-001"
            self.dimension=3072
        elif self.embedding_model_type=="models/text-embedding-004":
            self.embedding_model = "models/gemini-embedding-004"
            self.dimension=1536
        elif self.embedding_model_type=="fasttext":
            import fasttextembeding
            self.embedding_model = "fasttext"
            self.dimension=300
        if not os.path.exists("data"):
            os.makedirs("data")
        if not os.path.exists("data/vector_db"):
            os.makedirs("data/vector_db")
        if not os.path.exists(self.db_path):
            os.makedirs(self.db_path)
            
        if os.path.exists(self.index_file):
            self.index = faiss.read_index(self.index_file)
            with open(self.metadata_file, "rb") as f:
                self.metadata = pickle.load(f)
        else:
            self.index = faiss.IndexFlatL2(self.dimension)
            self.metadata = []

    def get_embedding(self, text):
        if self.embedding_model_type=="fasttext":
            import fasttextembeding
            return np.array(fasttextembeding.embedding_text(text)).astype('float32')
        else:
            # print(f"DEBUG: Getting embedding for: {text[:50]}...")
            response = self.client.models.embed_content(
                model=self.embedding_model,
                contents=[text]
            )
            # print("DEBUG: Embedding received.")
            return np.array(response.embeddings[0].values).astype('float32')

    def add_text(self, text, extra_metadata=None):
        if not text.strip():
            return
        
        embedding = self.get_embedding(text)
        embedding = np.expand_dims(embedding, axis=0)
        
        self.index.add(embedding)
        self.metadata.append({
            "at": datetime.datetime.now().timestamp(),
            "text": text,
            "extra": extra_metadata or {}
        })
        self._save()

    def add_texts(self, texts, extra_metadatas=None):
        if not texts:
            return
        
        print(f"DEBUG: Batch embedding {len(texts)} texts...")
        # Batch embedding for efficiency
        response = self.client.models.embed_content(
            model=self.embedding_model,
            contents=texts
        )
        print("DEBUG: Batch embeddings received.")
        embeddings = np.array([e.values for e in response.embeddings]).astype('float32')
        
        self.index.add(embeddings)
        for i, text in enumerate(texts):
            meta = {}
            if extra_metadatas and i < len(extra_metadatas):
                meta = extra_metadatas[i] or {}
            self.metadata.append({
                "at": datetime.datetime.now().timestamp(),
                "text": text,
                "extra": meta or {}
            })
        self._save()

    def search(self, query, top_k=5):
        if self.index.ntotal == 0:
            return []
            
        query_embedding = self.get_embedding(query)
        query_embedding = np.expand_dims(query_embedding, axis=0)
        
        distances, indices = self.index.search(query_embedding, top_k)
        
        results = []
        for i, idx in enumerate(indices[0]):
            if idx != -1 and idx < len(self.metadata):
                results.append({
                    "text": self.metadata[idx]["text"],
                    "extra": self.metadata[idx]["extra"],
                    "distance": float(distances[0][i])
                })
        return results

    def delete_texts(self, texts_to_delete):
        if not texts_to_delete:
            return
        
        if isinstance(texts_to_delete, str):
            texts_to_delete = [texts_to_delete]
            
        print(f"DEBUG: Deleting {len(texts_to_delete)} texts...")
        
        # Filter metadata and keep only non-matching entries
        new_metadata = []
        indices_to_keep = []
        
        for i, meta in enumerate(self.metadata):
            if meta["text"] not in texts_to_delete:
                new_metadata.append(meta)
                indices_to_keep.append(i)
        
        if len(new_metadata) == len(self.metadata):
            print("DEBUG: No matching texts found for deletion.")
            return

        if new_metadata:
            # Rebuild index with remaining vectors
            print(f"DEBUG: Rebuilding index with {len(new_metadata)} entries...")
            
            # Extract current vectors from existing index
            current_vectors = []
            for idx in indices_to_keep:
                vec = self.index.reconstruct(idx)
                current_vectors.append(vec)
            
            # Ensure the array is float32 and contiguous for FAISS
            embeddings = np.array(current_vectors).astype('float32')
            embeddings = np.ascontiguousarray(embeddings)
            
            new_index = faiss.IndexFlatL2(self.dimension)
            new_index.add(embeddings)
            self.index = new_index
        else:
            # If nothing left, just reset to empty index
            self.index = faiss.IndexFlatL2(self.dimension)
            
        self.metadata = new_metadata
        self._save()
        print("DEBUG: Deletion complete.")

    def _save(self):
        print(f"DEBUG: Saving index to {self.index_file}...")
        faiss.write_index(self.index, self.index_file)
        print(f"DEBUG: Saving metadata to {self.metadata_file}...")
        with open(self.metadata_file, "wb") as f:
            pickle.dump(self.metadata, f)
        print("DEBUG: Save complete.")

    def get_full_text(self, file_path, file_description=None):
        """
        Extract full text from a file based on its extension.
        If file_description is provided and not empty, it returns that.
        Otherwise, it reads the file content or calls Gemini for multimedia.
        """
        if file_description and file_description.strip():
            return( file_description, {"file_path":file_path})

        if not os.path.exists(file_path):
            print(f"ERROR: File not found: {file_path}")
            return ("","")

        ext = os.path.splitext(file_path)[1].lower()
        print(f"DEBUG: Extracting text from {file_path} (ext: {ext})...")

        try:
            if ext == ".txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    return (f.read(), {"file_path":file_path})

            elif ext == ".pdf":
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return (text, {"file_path":file_path})

            elif ext == ".docx":
                doc = Document(file_path)
                return ("\n".join([para.text for para in doc.paragraphs]), {"file_path":file_path})

            elif ext == ".pptx":
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return ("\n".join(text_runs), {"file_path":file_path})

            elif ext == ".json":
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    return (json.dumps(data, indent=2, ensure_ascii=False), {"file_path":file_path})

            elif ext in [".xlsx", ".xls", ".csv"]:
                if ext == ".csv":
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                return (df.to_string(index=False), {"file_path":file_path})

            elif ext in [".jpg", ".jpeg", ".png", ".webp", ".mp4", ".mov", ".avi", ".mp3", ".wav", ".m4a"]:
                print(f"DEBUG: Calling Gemini to describe multimedia file: {file_path}")
                
                # Upload file to Gemini for processing
                try:
                    # Get MIME type
                    import mimetypes
                    mime_type, _ = mimetypes.guess_type(file_path)
                    
                    with open(file_path, "rb") as f:
                        file_data = f.read()
                    
                    prompt = "Hãy mô tả nội dung của file này một cách chi tiết để sử dụng cho việc tìm kiếm RAG sau này. Nếu là văn bản trong ảnh/video, hãy trích xuất toàn bộ. Nếu là âm thanh, hãy chuyển soạn nội dung. Chỉ cần trả lại nội dung mô tả súc tích nhưng dầy đủ thông tin có thể có."
                    
                    response = self.client.models.generate_content(
                        model="gemini-2.0-flash", # Using flash for speed and multimodal support
                        contents=[
                            types.Content(
                                role="user",
                                parts=[
                                    types.Part.from_bytes(data=file_data, mime_type=mime_type),
                                    types.Part.from_text(text=prompt)
                                ]
                            )
                        ]
                    )
                    return (response.text,{"file_path":file_path,"response":response})
                except Exception as e:
                    print(f"ERROR during Gemini call: {e}")
                    return (f"Mô tả cho file {os.path.basename(file_path)}",{"file_path":file_path,"error":str(e)})

            else:
                # Default fallback for unknown extensions
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        return (f.read(),{"file_path":file_path,"response":"unknown extensions"})
                except Exception as e1:
                    return (f"Nội dung file: {os.path.basename(file_path)}",{"file_path":file_path,"error":str(e1)})

        except Exception as e:
            print(f"ERROR extracting text from {file_path}: {e}")
            return ("",{"file_path":file_path,"error":str(e)})

    def chunk_text(self, text, chunk_size=1000, overlap=200):
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start += chunk_size - overlap
        return chunks

    def add_file(self, file_path, file_description=None):
        """
        Index a file in the vector database.
        """
        text, metadata = self.get_full_text(file_path, file_description)
        print( f"==============> add_file: {text}")
        chunks = self.chunk_text(text)
        for chunk in chunks:
            self.add_text(chunk,metadata)
        self._save()


    def add_folder(self, folder_path):
        """
        Recursively index all supported files in a folder.
        """
        if not os.path.isdir(folder_path):
            print(f"ERROR: Folder not found or is not a directory: {folder_path}")
            return

        supported_extensions = [
            ".txt", ".pdf", ".docx", ".pptx", ".json", ".xlsx", ".xls", ".csv",
            ".jpg", ".jpeg", ".png", ".webp", ".mp4", ".mov", ".avi", ".mp3", ".wav", ".m4a"
        ]

        print(f"DEBUG: Scanning folder: {folder_path}...")
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in supported_extensions:
                    full_path = os.path.join(root, file)
                    print(f"DEBUG: Adding file from folder: {full_path}")
                    self.add_file(full_path)