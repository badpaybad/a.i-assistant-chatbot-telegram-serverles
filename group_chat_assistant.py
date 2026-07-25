import os
import re
import json
import sqlite3
import requests
from datetime import datetime
from typing import List, Dict, Any, Optional

# ==========================================
# THƯ VIỆN XỬ LÝ FILE (Office, Image, Audio)
# ==========================================
# 1. Document & Office
try:
    import docx  # python-docx
except ImportError:
    docx = None

try:
    import openpyxl  # openpyxl for Excel
except ImportError:
    openpyxl = None

try:
    import pptx  # python-pptx
except ImportError:
    pptx = None

try:
    import pypdf  # pypdf for PDF
except ImportError:
    pypdf = None

# 2. Image
try:
    from PIL import Image
    import pytesseract  # Tesseract OCR
except ImportError:
    Image = None
    pytesseract = None

# 3. Audio (Whisper local)
try:
    import whisper  # openai-whisper
except ImportError:
    whisper = None

# ==========================================
# CẤU HÌNH HỆ THỐNG
# ==========================================
OLLAMA_ENDPOINT = "http://localhost:11434/api/chat"
MODEL_NAME = "qwen2.5:7b"  # Khuyên dùng Qwen2.5, Llama3.1 hoặc Mistral-Nemo
DB_FILE = "bot_memory.db"
MAX_RETRIES = 2

# Cache cho Whisper Model để không nạp lại mỗi lần đọc audio
WHISPER_MODEL = None

# ==========================================
# 1. KHỞI TẠO DATABASE (SQLITE LOCAL)
# ==========================================
def init_db():
    """Khởi tạo SQLite lưu trữ Bộ nhớ dài hạn và Nhắc nhở."""
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS memory (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            category TEXT,
            content TEXT NOT NULL,
            created_by TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS reminders (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            assignee TEXT,
            remind_at TEXT NOT NULL,
            status TEXT DEFAULT 'pending',
            created_by TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    conn.commit()
    conn.close()

init_db()

# ==========================================
# 2. CÁC HÀM XỬ LÝ FILE CHUYÊN SÂU (OFFICE, IMAGE, AUDIO)
# ==========================================

def tool_read_document_file(file_path: str) -> str:
    """Đọc và trích xuất nội dung từ các file văn bản và tài liệu Office (.docx, .xlsx, .pptx, .pdf, .txt, .csv, .json)."""
    if not os.path.exists(file_path):
        return f"[Lỗi: Không tìm thấy file tại {file_path}]"
    
    ext = os.path.splitext(file_path)[1].lower()
    try:
        # 1. File Word (.docx)
        if ext == ".docx":
            if docx is None:
                return "[Lỗi: Thiếu thư viện 'python-docx'. Chạy: pip install python-docx]"
            doc = docx.Document(file_path)
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(full_text)[:3000]

        # 2. File Excel (.xlsx)
        elif ext in [".xlsx", ".xls"]:
            if openpyxl is None:
                return "[Lỗi: Thiếu thư viện 'openpyxl'. Chạy: pip install openpyxl]"
            wb = openpyxl.load_workbook(file_path, data_only=True)
            output = []
            for sheet_name in wb.sheetnames[:3]: # Đọc tối đa 3 sheet đầu
                sheet = wb[sheet_name]
                output.append(f"--- Sheet: {sheet_name} ---")
                for row in list(sheet.iter_rows(values_only=True))[:30]: # Đọc 30 dòng đầu
                    row_vals = [str(cell) for cell in row if cell is not None]
                    if row_vals:
                        output.append(" | ".join(row_vals))
            return "\n".join(output)[:3000]

        # 3. File PowerPoint (.pptx)
        elif ext == ".pptx":
            if pptx is None:
                return "[Lỗi: Thiếu thư viện 'python-pptx'. Chạy: pip install python-pptx]"
            prs = pptx.Presentation(file_path)
            slides_text = []
            for idx, slide in enumerate(prs.slides[:10]): # Max 10 slides
                slides_text.append(f"--- Slide {idx+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides_text.append(shape.text.strip())
            return "\n".join(slides_text)[:3000]

        # 4. File PDF (.pdf)
        elif ext == ".pdf":
            if pypdf is None:
                return "[Lỗi: Thiếu thư viện 'pypdf'. Chạy: pip install pypdf]"
            reader = pypdf.PdfReader(file_path)
            text = ""
            for page in reader.pages[:5]: # Max 5 trang
                text += page.extract_text() or ""
            return text[:3000]

        # 5. File Text thuần (.txt, .md, .csv, .json, .log)
        elif ext in [".txt", ".md", ".json", ".csv", ".log", ".py"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:3000]

        else:
            return f"[Lỗi: Định dạng tài liệu '{ext}' chưa được hỗ trợ]"

    except Exception as e:
        return f"[Lỗi bóc tách tài liệu: {str(e)}]"


def tool_process_image_file(file_path: str) -> str:
    """Trích xuất chữ (OCR) và đọc thông tin chi tiết từ file hình ảnh (.png, .jpg, .jpeg, .webp)."""
    if not os.path.exists(file_path):
        return f"[Lỗi: Không tìm thấy ảnh tại {file_path}]"
    
    if Image is None:
        return "[Lỗi: Thiếu thư viện Pillow. Chạy: pip install Pillow pytesseract]"

    try:
        img = Image.open(file_path)
        width, height = img.size
        format_name = img.format
        info_str = f"Chỉ số ảnh: {format_name}, Kích thước: {width}x{height}px.\n"

        # Trích xuất văn bản qua Tesseract OCR nếu có cài đặt
        ocr_text = ""
        if pytesseract is not None:
            try:
                ocr_text = pytesseract.image_to_string(img, lang="vie+eng").strip()
            except Exception:
                ocr_text = "(Chưa cấu hình binary Tesseract OCR trên OS)"
        
        if ocr_text:
            return f"{info_str}Văn bản trích xuất từ ảnh (OCR):\n{ocr_text[:2000]}"
        else:
            return f"{info_str}(Không tìm thấy chữ trong ảnh hoặc chưa cài Tesseract engine)"

    except Exception as e:
        return f"[Lỗi xử lý hình ảnh: {str(e)}]"


def tool_transcribe_audio_file(file_path: str) -> str:
    """Chuyển đổi file ghi âm/âm thanh (.mp3, .wav, .m4a, .ogg) thành văn bản tiếng Việt/Anh bằng Whisper."""
    global WHISPER_MODEL
    if not os.path.exists(file_path):
        return f"[Lỗi: Không tìm thấy file audio tại {file_path}]"
    
    if whisper is None:
        return "[Lỗi: Thiếu thư viện 'openai-whisper'. Chạy: pip install openai-whisper ffmpeg-python]"

    try:
        # Load model Whisper (Base/Small tối ưu trên GPU local)
        if WHISPER_MODEL is None:
            print("🎙️ [SYSTEM]: Đang nạp mô hình Whisper local...")
            WHISPER_MODEL = whisper.load_model("base")
            
        result = WHISPER_MODEL.transcribe(file_path)
        transcribed_text = result.get("text", "").strip()
        return f"Nội dung giọng nói bóc tách từ Audio:\n\"{transcribed_text}\""

    except Exception as e:
        return f"[Lỗi khi bóc tách audio Whisper: {str(e)}]"


def tool_save_memory(content: str, category: str = "general", created_by: str = "System") -> str:
    """Lưu thông tin quan trọng, quy định, quyết định vào DB bộ nhớ dài hạn."""
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()
    cursor.execute("INSERT INTO memory (category, content, created_by) VALUES (?, ?, ?)", (category, content, created_by))
    conn.commit()
    conn.close()
    return f"✅ Đã lưu vào bộ nhớ DB: '{content}' (Phân loại: {category})"


def tool_search_memory(query: str) -> str:
    """Tra cứu các quy định, thông tin đã lưu trong bộ nhớ DB."""
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()
    cursor.execute("SELECT category, content, created_by, created_at FROM memory WHERE content LIKE ?", (f"%{query}%",))
    rows = cursor.fetchall()
    conn.close()
    if not rows:
        return "Không tìm thấy thông tin phù hợp trong DB bộ nhớ."
    results = [f"- [{r[0].upper()}] {r[1]} (Tạo bởi: {r[2]} ngày {r[3]})" for r in rows]
    return "Kết quả từ DB bộ nhớ:\n" + "\n".join(results)


def tool_create_reminder(title: str, remind_at: str, assignee: str = "Mọi người", created_by: str = "System") -> str:
    """Tạo lịch nhắc nhở/giao việc cho thành viên trong nhóm."""
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()
    cursor.execute("INSERT INTO reminders (title, assignee, remind_at, created_by) VALUES (?, ?, ?, ?)", (title, assignee, remind_at, created_by))
    conn.commit()
    conn.close()
    return f"⏰ Đã tạo nhắc nhở: '{title}' cho [{assignee}] vào lúc [{remind_at}]."

# ==========================================
# 3. SCHEMA DÀNH CHO LLM TOOL CALLING
# ==========================================
TOOLS_SCHEMA = [
    {
        "type": "function",
        "function": {
            "name": "read_document_file",
            "description": "Đọc nội dung văn bản từ các file tài liệu Office (.docx, .xlsx, .pptx, .pdf, .txt, .csv, .json) được đính kèm.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Đường dẫn file tài liệu."}
                },
                "required": ["file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "process_image_file",
            "description": "Bóc tách chữ (OCR) và xem thuộc tính từ file hình ảnh (.png, .jpg, .jpeg, .webp).",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Đường dẫn file ảnh."}
                },
                "required": ["file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "transcribe_audio_file",
            "description": "Chuyển file ghi âm giọng nói (.mp3, .wav, .m4a, .ogg) thành văn bản.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Đường dẫn file âm thanh."}
                },
                "required": ["file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "save_memory",
            "description": "Lưu lại quy định, quyết định, thông tin quan trọng của nhóm vào DB.",
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {"type": "string", "description": "Nội dung cần nhớ."},
                    "category": {"type": "string", "description": "Phân loại (quy_trinh, vpn, server, note)."}
                },
                "required": ["content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_memory",
            "description": "Tra cứu các thông tin, quy trình đã lưu trong DB.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Từ khóa tìm kiếm."}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_reminder",
            "description": "Tạo nhắc nhở hẹn giờ hoặc giao task cho thành viên.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Tên công việc."},
                    "remind_at": {"type": "string", "description": "Thời gian nhắc."},
                    "assignee": {"type": "string", "description": "Người nhận task."}
                },
                "required": ["title", "remind_at"]
            }
        }
    }
]

AVAILABLE_FUNCTIONS = {
    "read_document_file": tool_read_document_file,
    "process_image_file": tool_process_image_file,
    "transcribe_audio_file": tool_transcribe_audio_file,
    "save_memory": tool_save_memory,
    "search_memory": tool_search_memory,
    "create_reminder": tool_create_reminder
}

# ==========================================
# 4. PIPELINE ĐIỀU PHỐI (ORCHESTRATOR)
# ==========================================
def build_system_prompt() -> str:
    now_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    return f"""Bạn là Trợ lý Công việc (Work Assistant) trong nhóm chat nội bộ.
Thời gian hiện tại: {now_str}

[QUY TẮC BẮT BUỘC GỌI TOOL]
1. Nếu user gửi hoặc yêu cầu đọc tài liệu (.docx, .xlsx, .pptx, .pdf, .txt, .csv) -> Gọi `read_document_file`.
2. Nếu user gửi file hình ảnh (.png, .jpg, .webp) -> Gọi `process_image_file`.
3. Nếu user gửi file âm thanh/ghi âm (.mp3, .wav, .m4a) -> Gọi `transcribe_audio_file`.
4. Nếu user muốn ghi nhớ/lưu thông tin -> Gọi `save_memory`.
5. Nếu user hỏi quy trình/thông tin cũ -> Gọi `search_memory`.
6. Nếu user nhờ nhắc lịch/giao việc -> Gọi `create_reminder`.
7. Trả lời luôn ngắn gọn, vào trọng tâm và lịch sự."""


def run_pipeline(
    speaker_name: str,
    speaker_role: str,
    user_query: str,
    raw_chat_history: List[Dict[str, str]],
    attached_file_path: Optional[str] = None
) -> str:
    """Hàm chạy chính điều hướng Tool Calls và trả về câu trả lời."""
    
    # 1. Tóm tắt lịch sử chat
    history_str = "\n".join([f"[{m.get('timestamp','')}] {m.get('sender','')}: {m.get('text','')}" for m in raw_chat_history[-8:]])
    
    file_info = f"\n<attached_file_path>{attached_file_path}</attached_file_path>" if attached_file_path else ""

    user_content = f"""<group_context>
{history_str}
</group_context>

<current_speaker>
Tên: {speaker_name} ({speaker_role})
</current_speaker>{file_info}

<question>
{user_query}
</question>"""

    messages = [
        {"role": "system", "content": build_system_prompt()},
        {"role": "user", "content": user_content}
    ]

    payload = {
        "model": MODEL_NAME,
        "messages": messages,
        "tools": TOOLS_SCHEMA,
        "stream": False,
        "options": {"temperature": 0.1, "num_ctx": 4096}
    }

    try:
        # Lượt 1: LLM quyết định Tool Call
        response = requests.post(OLLAMA_ENDPOINT, json=payload, timeout=35).json()
        message = response.get("message", {})

        # Nếu có Tool Call
        if message.get("tool_calls"):
            messages.append(message)
            
            for tool_call in message["tool_calls"]:
                fn_name = tool_call["function"]["name"]
                fn_args = tool_call["function"]["arguments"]

                if fn_name in ["save_memory", "create_reminder"]:
                    fn_args["created_by"] = speaker_name

                print(f"🛠️ [BOT EXECUTE TOOL]: {fn_name}({fn_args})")

                if fn_name in AVAILABLE_FUNCTIONS:
                    tool_output = AVAILABLE_FUNCTIONS[fn_name](**fn_args)
                else:
                    tool_output = f"[Lỗi: Tool {fn_name} không tồn tại]"

                messages.append({
                    "role": "tool",
                    "content": str(tool_output)
                })

            # Lượt 2: LLM tổng hợp kết quả từ Tool
            final_payload = {
                "model": MODEL_NAME,
                "messages": messages,
                "stream": False,
                "options": {"temperature": 0.1}
            }
            final_res = requests.post(OLLAMA_ENDPOINT, json=final_payload, timeout=35).json()
            return final_res.get("message", {}).get("content", "").strip()

        else:
            return message.get("content", "").strip()

    except Exception as e:
        return f"[Lỗi xử lý Pipeline: {str(e)}]"

# ==========================================
# 5. DEMO CHẠY THỬ
# ==========================================
if __name__ == "__main__":
    chat_history = [{"timestamp": "08:00", "sender": "Nam (Lead)", "text": "Chào cả nhà, chúc tuần mới vui vẻ."}]

    # Demo 1: Đọc file Excel
    print("=== DEMO 1: ĐỌC FILE EXCEL BẢNG LƯƠNG/CHI PHÍ ===")
    out_excel = run_pipeline(
        speaker_name="Minh (Dev)",
        speaker_role="Engineer",
        user_query="@Bot đọc giúp mình file Excel này xem tổng chi phí là bao nhiêu?",
        raw_chat_history=chat_history,
        attached_file_path="danh_sach_chi_phi.xlsx"
    )
    print("Bot Reply:", out_excel)
    print("-" * 50)

    # Demo 2: Đọc file Word (.docx)
    print("\n=== DEMO 2: ĐỌC FILE WORD BÁO CÁO ===")
    out_word = run_pipeline(
        speaker_name="Hoa (QA)",
        speaker_role="Tester",
        user_query="@Bot tóm tắt ngắn gọn file Word hợp đồng này giúp chị?",
        raw_chat_history=chat_history,
        attached_file_path="hop_dong_dich_vu.docx"
    )
    print("Bot Reply:", out_word)
    print("-" * 50)

    # Demo 3: Xử lý Voice Audio bằng Whisper
    print("\n=== DEMO 3: BÓC TÁCH FILE VOICE AUDIO (.MP3/.WAV) ===")
    out_audio = run_pipeline(
        speaker_name="Nam (Lead)",
        speaker_role="Leader",
        user_query="@Bot nghe giúp anh đoạn voice ghi âm cuộc họp này xem sếp dặn gì?",
        raw_chat_history=chat_history,
        attached_file_path="voice_note_cuoc_hop.mp3"
    )
    print("Bot Reply:", out_audio)